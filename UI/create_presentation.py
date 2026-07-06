import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Palette definition (Slate & Emerald theme)
    NAVY = RGBColor(15, 23, 42)          # Slate 900 (#0f172a)
    EMERALD = RGBColor(16, 185, 129)     # Emerald 500 (#10b981)
    CHARCOAL = RGBColor(51, 65, 85)      # Slate 700 (#334155)
    SLATE_400 = RGBColor(148, 163, 184)  # Slate 400 (#94a3b8)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50 (#f8fafc)
    CARD_BG = RGBColor(239, 246, 255)    # Blue 50 (#eff6ff)
    BORDER_COLOR = RGBColor(219, 234, 254) # Blue 100 (#dbeafe)
    
    def set_solid_fill(shape, color):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        
    def add_background(slide, color):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        set_solid_fill(rect, color)
        rect.line.fill.background() # No border line

    def add_header(slide, title_text, category_text="EXAMAI: AI-POWERED SECURE EXAMINATION PORTAL"):
        # Thin top border line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
        set_solid_fill(line, EMERALD)
        line.line.fill.background()
        
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.733), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Trebuchet MS'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = EMERALD
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_top = Inches(0)
        tf_title.margin_bottom = Inches(0)
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Georgia'
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY
        
    def add_bullet_point(tf, bold_prefix, text_body, size=16, space_before=14):
        p = tf.add_paragraph() if (tf.paragraphs and tf.paragraphs[0].text) else tf.paragraphs[0]
        p.space_before = Pt(space_before)
        p.level = 0
        
        # Add bold prefix
        r1 = p.add_run()
        r1.text = bold_prefix
        r1.font.name = 'Trebuchet MS'
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = NAVY
        
        # Add normal body
        r2 = p.add_run()
        r2.text = text_body
        r2.font.name = 'Trebuchet MS'
        r2.font.size = Pt(size)
        r2.font.color.rgb = CHARCOAL
        
    def add_card(slide, title, items, x, y, width, height):
        # Background card shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
        set_solid_fill(card, CARD_BG)
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        # Text frame inside card
        tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.25), Inches(width - 0.5), Inches(height - 0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        
        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = "•  " + item
            p_item.font.name = 'Trebuchet MS'
            p_item.font.size = Pt(12)
            p_item.font.color.rgb = CHARCOAL
            p_item.space_before = Pt(6)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # ----------------------------------------------------
    slide_1 = prs.slides.add_slide(blank_layout)
    add_background(slide_1, NAVY)
    
    # Top Accent Border
    accent_bar = slide_1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    set_solid_fill(accent_bar, EMERALD)
    accent_bar.line.fill.background()
    
    # College Name Box
    tx_college = slide_1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(0.8))
    tf_col = tx_college.text_frame
    tf_col.word_wrap = True
    p_col = tf_col.paragraphs[0]
    p_col.text = "[YOUR COLLEGE/UNIVERSITY NAME]"
    p_col.font.name = 'Trebuchet MS'
    p_col.font.size = Pt(20)
    p_col.font.bold = True
    p_col.font.color.rgb = EMERALD
    
    # Project Title Box
    tx_title = slide_1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(2.2))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "ExamAI"
    p_title.font.name = 'Georgia'
    p_title.font.size = Pt(54)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    
    p_sub = tf_title.add_paragraph()
    p_sub.text = "AI-Powered Secure Examination Portal"
    p_sub.font.name = 'Trebuchet MS'
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = EMERALD
    p_sub.space_before = Pt(12)
    
    # Presenter Box
    tx_pres = slide_1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(6.0), Inches(1.8))
    tf_pres = tx_pres.text_frame
    tf_pres.word_wrap = True
    
    p_lbl = tf_pres.paragraphs[0]
    p_lbl.text = "Presented By:"
    p_lbl.font.name = 'Trebuchet MS'
    p_lbl.font.size = Pt(14)
    p_lbl.font.color.rgb = EMERALD
    
    p_name = tf_pres.add_paragraph()
    p_name.text = "Sumangal Kayal"
    p_name.font.name = 'Trebuchet MS'
    p_name.font.size = Pt(24)
    p_name.font.bold = True
    p_name.font.color.rgb = WHITE
    p_name.space_before = Pt(4)
    
    p_det = tf_pres.add_paragraph()
    p_det.text = "B.Tech Final Year Project\nDepartment of Computer Science & Engineering"
    p_det.font.name = 'Trebuchet MS'
    p_det.font.size = Pt(13)
    p_det.font.color.rgb = SLATE_400
    p_det.space_before = Pt(4)
    
    # Side Decoration Box
    side_dec = slide_1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(2.2), Inches(2.8), Inches(3.5))
    set_solid_fill(side_dec, NAVY)
    side_dec.line.color.rgb = EMERALD
    side_dec.line.width = Pt(2)
    tf_dec = side_dec.text_frame
    tf_dec.word_wrap = True
    tf_dec.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_dec = tf_dec.paragraphs[0]
    p_dec.text = "B.TECH CAPSTONE\n\nEXAMINATION\nPORTAL\n\nAI PROCTORING"
    p_dec.font.name = 'Trebuchet MS'
    p_dec.font.size = Pt(16)
    p_dec.font.bold = True
    p_dec.font.color.rgb = EMERALD
    p_dec.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 2: Index / Agenda
    # ----------------------------------------------------
    slide_2 = prs.slides.add_slide(blank_layout)
    add_background(slide_2, LIGHT_BG)
    add_header(slide_2, "Index & Presentation Outline")
    
    # Left Agenda Column
    tx_ind1 = slide_2.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
    tf_ind1 = tx_ind1.text_frame
    tf_ind1.word_wrap = True
    
    add_bullet_point(tf_ind1, "01.  ", "Project Overview & Introduction", size=18, space_before=15)
    add_bullet_point(tf_ind1, "02.  ", "Problem Statement & Project Scope", size=18, space_before=15)
    add_bullet_point(tf_ind1, "03.  ", "Core System Features (Tri-Role Panel)", size=18, space_before=15)
    add_bullet_point(tf_ind1, "04.  ", "System Architecture & Flow", size=18, space_before=15)
    add_bullet_point(tf_ind1, "05.  ", "AI Proctoring & Verification (face-api.js)", size=18, space_before=15)
    
    # Right Agenda Column
    tx_ind2 = slide_2.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.2), Inches(5.0))
    tf_ind2 = tx_ind2.text_frame
    tf_ind2.word_wrap = True
    
    add_bullet_point(tf_ind2, "06.  ", "Database Schema & API Design", size=18, space_before=15)
    add_bullet_point(tf_ind2, "07.  ", "Technical Implementation Stack", size=18, space_before=15)
    add_bullet_point(tf_ind2, "08.  ", "System Workflow & Process Steps", size=18, space_before=15)
    add_bullet_point(tf_ind2, "09.  ", "Future Enhancements & Additions", size=18, space_before=15)
    add_bullet_point(tf_ind2, "10.  ", "Academic Relevance & Final Summary", size=18, space_before=15)

    # ----------------------------------------------------
    # SLIDE 3: Project Overview & Introduction
    # ----------------------------------------------------
    slide_3 = prs.slides.add_slide(blank_layout)
    add_background(slide_3, LIGHT_BG)
    add_header(slide_3, "Introduction & Project Overview")
    
    tx_3 = slide_3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.0))
    tf_3 = tx_3.text_frame
    tf_3.word_wrap = True
    
    add_bullet_point(tf_3, "Next-Gen Assessment Portal: ", "ExamAI is a modern, full-stack online examination system designed to manage remote assessments securely.")
    add_bullet_point(tf_3, "Roles & Capabilities: ", "Integrates three distinct functional roles - Student (examinee), Examiner (invigilator/creator), and Administrator (system audit/oversight).")
    add_bullet_point(tf_3, "Intelligent Monitoring: ", "Protects academic integrity using client-side AI analysis to capture anomalies (window tab switches, facial absences, multiple users).")
    add_bullet_point(tf_3, "Autonomous Pipeline: ", "Streamlines assessment creation, answer evaluation, risk log processing, and final dashboard analytics in a single environment.")
    
    add_card(slide_3, "ExamAI Highlights", 
             ["MERN Stack + Client AI", "Tri-Role Access Boundaries", "Real-Time Incident Tracking", "Automated Evaluation Engine"],
             8.5, 1.5, 4.0, 5.0)

    # ----------------------------------------------------
    # SLIDE 4: Problem Statement
    # ----------------------------------------------------
    slide_4 = prs.slides.add_slide(blank_layout)
    add_background(slide_4, LIGHT_BG)
    add_header(slide_4, "Problem Statement & Necessity")
    
    tx_4 = slide_4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.0))
    tf_4 = tx_4.text_frame
    tf_4.word_wrap = True
    
    add_bullet_point(tf_4, "Integrity Risks: ", "Traditional remote tests lack direct oversight, leading to widespread cheating, tab switching, and unauthorized aids.")
    add_bullet_point(tf_4, "Scalability Limits: ", "Manual remote invigilation requires a high ratio of human proctors, causing excessive costs and human oversight fatigue.")
    add_bullet_point(tf_4, "LMS Shortcomings: ", "Standard Learning Management Systems do not include integrated facial authentication or active background behavioral tracking.")
    add_bullet_point(tf_4, "Unified Platform Need: ", "Institutions need a cost-effective, real-time proctored test environment that evaluates students and tracks risk indexes automatically.")
    
    add_card(slide_4, "Core Target Problems", 
             ["Unverified candidate identity", "Frequent browser tab switches", "Exorbitant cost of human invigilation", "Scattered audit logs & charts"],
             8.5, 1.5, 4.0, 5.0)

    # ----------------------------------------------------
    # SLIDE 5: Core System Features
    # ----------------------------------------------------
    slide_5 = prs.slides.add_slide(blank_layout)
    add_background(slide_5, LIGHT_BG)
    add_header(slide_5, "Core System Features")
    
    tx_5 = slide_5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.0))
    tf_5 = tx_5.text_frame
    tf_5.word_wrap = True
    
    add_bullet_point(tf_5, "Student Dashboard: ", "Features Secure login, Face-ID validation, active exam panel, automatic countdown autosave, and instantaneous grade delivery.")
    add_bullet_point(tf_5, "Examiner Workspace: ", "Provides interface to design exams, author custom questions, review submissions, and inspect real-time anomaly timelines with risk scoring.")
    add_bullet_point(tf_5, "AI Proctoring Module: ", "Tracks tab switches, window minimizations, facial presence, and calculates composite candidate risk quotients.")
    
    add_card(slide_5, "Admin Controls", 
             ["System-wide analytics dashboard", "Interactive audit log viewer", "Database model security controls", "Account creation & management"],
             8.5, 1.5, 4.0, 5.0)

    # ----------------------------------------------------
    # SLIDE 6: System Architecture
    # ----------------------------------------------------
    slide_6 = prs.slides.add_slide(blank_layout)
    add_background(slide_6, LIGHT_BG)
    add_header(slide_6, "System Architecture & Flow")
    
    tx_6 = slide_6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.0))
    tf_6 = tx_6.text_frame
    tf_6.word_wrap = True
    
    add_bullet_point(tf_6, "Client-Server Decoupling: ", "React frontend communicates with Express backend through structured REST API endpoints.")
    add_bullet_point(tf_6, "Security Architecture: ", "Stateless authentication via JSON Web Tokens (JWT) stored client-side; password encryption utilizing bcryptjs.")
    add_bullet_point(tf_6, "Database Management: ", "Express communicates with MongoDB Atlas via Mongoose ODM for reliable schemas.")
    add_bullet_point(tf_6, "AI Pipeline: ", "Face detection model runs fully on the client browser via face-api.js, saving server resources and bandwidth.")
    
    add_card(slide_6, "Data Flow Steps", 
             ["1. Auth & JWT generation", "2. Client-side face descriptor check", "3. Exam delivery & active timers", "4. Event logs pushed to server DB", "5. Live examiner dashboard updates"],
             8.5, 1.5, 4.0, 5.0)

    # ----------------------------------------------------
    # SLIDE 7: AI Proctoring & Identity Verification
    # ----------------------------------------------------
    slide_7 = prs.slides.add_slide(blank_layout)
    add_background(slide_7, LIGHT_BG)
    add_header(slide_7, "AI Proctoring & Face-ID Verification")
    
    tx_7 = slide_7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.0))
    tf_7 = tx_7.text_frame
    tf_7.word_wrap = True
    
    add_bullet_point(tf_7, "Face Descriptors: ", "Extracts a 128-dimensional mathematical vector (descriptor) representing distinct facial features via deep learning models.")
    add_bullet_point(tf_7, "Verification Process: ", "During verification, face-api.js matches live webcam descriptors against the registered vector using Euclidean distance.")
    add_bullet_point(tf_7, "Tab Switch Logging: ", "Intercepts browser Page Visibility API events to record when a candidate leaves the examination window.")
    add_bullet_point(tf_7, "Risk Scoring: ", "Computes a cumulative danger metric based on incident logs. Examiners get immediate telemetry on candidate violations.")
    
    add_card(slide_7, "Monitored Violations", 
             ["Tab changes & focus losses", "No face detected in webcam", "Identity validation failures", "Audio/visual threshold warnings"],
             8.5, 1.5, 4.0, 5.0)

    # ----------------------------------------------------
    # SLIDE 8: Database & API Design
    # ----------------------------------------------------
    slide_8 = prs.slides.add_slide(blank_layout)
    add_background(slide_8, LIGHT_BG)
    add_header(slide_8, "Database Models & API Modules")
    
    tx_8 = slide_8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.0))
    tf_8 = tx_8.text_frame
    tf_8.word_wrap = True
    
    add_bullet_point(tf_8, "Mongoose Schema Models: ", "Utilizes structured models for Users (roles, hashes, face descriptors), Exams, Submissions, and Incidents.")
    add_bullet_point(tf_8, "Auth Modules: ", "Supports registration, JWT token generation, and secure payload decryption on middleware routes.")
    add_bullet_point(tf_8, "Exam & Submission APIs: ", "Facilitates CRUD operations on questions, real-time response commits, and automated score compiling on exam termination.")
    add_bullet_point(tf_8, "Incident Logging APIs: ", "Stores timelines of violations and fetches analytics datasets for examiners.")
    
    add_card(slide_8, "Core API Endpoints", 
             ["POST /api/auth/login", "POST /api/auth/register-face", "GET/POST /api/exams", "POST /api/submissions", "POST /api/incidents/log"],
             8.5, 1.5, 4.0, 5.0)

    # ----------------------------------------------------
    # SLIDE 9: Technical Stack
    # ----------------------------------------------------
    slide_9 = prs.slides.add_slide(blank_layout)
    add_background(slide_9, LIGHT_BG)
    add_header(slide_9, "Technical Stack Details")
    
    tx_9 = slide_9.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.0))
    tf_9 = tx_9.text_frame
    tf_9.word_wrap = True
    
    add_bullet_point(tf_9, "Frontend (React 19 & Vite 8): ", "Responsive UI utilizing modern react hooks, React Router 7 for secure routing boundaries, and vanilla CSS for styling.")
    add_bullet_point(tf_9, "Backend (Node.js & Express 5): ", "Robust, async API layer, CORS support, secure file processing via Multer, and routing middlewares.")
    add_bullet_point(tf_9, "Database (MongoDB Atlas): ", "Flexible document-based storage, high availability indexing, and cloud replication.")
    add_bullet_point(tf_9, "Telemetry & AI: ", "Recharts for telemetry dashboards, Lucide React icons, and local face-api.js web models.")
    
    add_card(slide_9, "Key Libraries", 
             ["face-api.js (Local ML)", "Recharts (Visual analytics)", "bcryptjs (Password hash)", "jsonwebtoken (Auth token)"],
             8.5, 1.5, 4.0, 5.0)

    # ----------------------------------------------------
    # SLIDE 10: Future Scope & Enhancements
    # ----------------------------------------------------
    slide_10 = prs.slides.add_slide(blank_layout)
    add_background(slide_10, LIGHT_BG)
    add_header(slide_10, "Future Scope & Enhancements")
    
    tx_10 = slide_10.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.0))
    tf_10 = tx_10.text_frame
    tf_10.word_wrap = True
    
    add_bullet_point(tf_10, "Live Audio Monitoring: ", "Integration of noise capture algorithms to flag verbal discussions and background help anomalies.")
    add_bullet_point(tf_10, "Dedicated Lockdown Browser: ", "Building an OS-specific wrapper application to disable screenshots, secondary displays, and shortcut commands.")
    add_bullet_point(tf_10, "Server-Side AI Audits: ", "Moving heavy visual models to backend nodes for periodic feed reviews.")
    add_bullet_point(tf_10, "Automated Notification System: ", "Incorporating SMS or email auto-warnings when a student violates proctor bounds multiple times.")
    
    add_card(slide_10, "Value Additions", 
             ["Multi-face detection models", "Automatic exam certificate generation", "Live invigilator audio/video calls", "Predictive ML cheating models"],
             8.5, 1.5, 4.0, 5.0)

    # ----------------------------------------------------
    # SLIDE 11: Academic Relevance & Conclusion
    # ----------------------------------------------------
    slide_11 = prs.slides.add_slide(blank_layout)
    add_background(slide_11, LIGHT_BG)
    add_header(slide_11, "Academic Relevance & Summary")
    
    tx_11 = slide_11.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.0))
    tf_11 = tx_11.text_frame
    tf_11.word_wrap = True
    
    add_bullet_point(tf_11, "MERN Demonstration: ", "Showcases professional full-stack Javascript mastery, data schema optimizations, and security controls.")
    add_bullet_point(tf_11, "Browser ML Implementation: ", "Illustrates the capability of executing AI models directly in-browser for zero-cost client processing.")
    add_bullet_point(tf_11, "Enterprise Readiness: ", "Maintains strict separation of concerns, JWT stateless routes, and interactive admin dashboards.")
    add_bullet_point(tf_11, "Conclusion: ", "ExamAI resolves core online testing vulnerabilities, proving that smart software solutions can effectively secure distant learning.")
    
    add_card(slide_11, "Presentation Summary", 
             ["Secure MERN Architecture", "Real-Time Client-Side AI", "Optimized DB schemas", "Flexible Future-Ready Design"],
             8.5, 1.5, 4.0, 5.0)

    # ----------------------------------------------------
    # SLIDE 12: Thank You Slide (Dark Theme)
    # ----------------------------------------------------
    slide_12 = prs.slides.add_slide(blank_layout)
    add_background(slide_12, NAVY)
    
    # Top Accent Border
    accent_bar2 = slide_12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    set_solid_fill(accent_bar2, EMERALD)
    accent_bar2.line.fill.background()
    
    # Large THANK YOU text
    tx_ty = slide_12.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
    tf_ty = tx_ty.text_frame
    tf_ty.word_wrap = True
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "THANK YOU!"
    p_ty.font.name = 'Georgia'
    p_ty.font.size = Pt(64)
    p_ty.font.bold = True
    p_ty.font.color.rgb = WHITE
    p_ty.alignment = PP_ALIGN.LEFT
    
    p_ty_sub = tf_ty.add_paragraph()
    p_ty_sub.text = "Questions & Discussions are Welcome."
    p_ty_sub.font.name = 'Trebuchet MS'
    p_ty_sub.font.size = Pt(22)
    p_ty_sub.font.color.rgb = EMERALD
    p_ty_sub.space_before = Pt(12)
    
    # Details Box
    tx_ty_det = slide_12.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(8.0), Inches(1.8))
    tf_ty_det = tx_ty_det.text_frame
    tf_ty_det.word_wrap = True
    p_ty_det = tf_ty_det.paragraphs[0]
    p_ty_det.text = "Project: ExamAI - Secure Examination Portal"
    p_ty_det.font.name = 'Trebuchet MS'
    p_ty_det.font.size = Pt(16)
    p_ty_det.font.bold = True
    p_ty_det.font.color.rgb = WHITE
    
    p_ty_name = tf_ty_det.add_paragraph()
    p_ty_name.text = "Student Name: Sumangal Kayal | B.Tech CSE Final Year"
    p_ty_name.font.name = 'Trebuchet MS'
    p_ty_name.font.size = Pt(14)
    p_ty_name.font.color.rgb = SLATE_400
    p_ty_name.space_before = Pt(6)
    
    p_ty_col = tf_ty_det.add_paragraph()
    p_ty_col.text = "College Name: [Your College/University Name]"
    p_ty_col.font.name = 'Trebuchet MS'
    p_ty_col.font.size = Pt(14)
    p_ty_col.font.color.rgb = SLATE_400
    p_ty_col.space_before = Pt(4)

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "ExamAI_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    build_deck()
